import fitz
import re
import os
import json
import io
from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration
from typing import List, Dict

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

SUPPORTED_EXTENSIONS = {'.pdf', '.pptx', '.docx', '.txt', '.md', '.csv'}


class DocumentProcessor:

    def __init__(self):
        self.processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        self.model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

    def process_file(self, file_path: str) -> List[Dict]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return self.process_pdf(file_path)
        elif ext == '.pptx':
            return self.process_pptx(file_path)
        elif ext == '.docx':
            return self.process_docx(file_path)
        elif ext in ('.txt', '.md', '.csv'):
            return self.process_text_file(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")

    def process_pdf(self, pdf_path: str) -> List[Dict]:
        doc = fitz.open(pdf_path)
        body_font_size = self._analyze_body_font_size(doc)

        chunks = self._get_structure_based_chunks(doc, pdf_path, body_font_size)

        if len(chunks) < 3:
            chunks = self._get_paragraph_chunks(doc, pdf_path)

        doc.close()
        return chunks

    def _analyze_body_font_size(self, doc) -> float:
        font_counts = {}
        for page in doc:
            blocks = page.get_text("dict")["blocks"]
            for b in blocks:
                if "lines" in b:
                    for line in b["lines"]:
                        for span in line["spans"]:
                            size = round(span["size"], 1)
                            text = span["text"].strip()
                            if len(text) > 1:
                                font_counts[size] = font_counts.get(size, 0) + len(text)

        if not font_counts:
            return 11.0
        return max(font_counts, key=font_counts.get)

    def _clean_text(self, text: str) -> str:
        return re.sub(r'\s+', ' ', text).strip()

    def _generate_caption(self, image_bytes: bytes) -> str:
        try:
            image = Image.open(io.BytesIO(image_bytes)).convert('RGB')
            inputs = self.processor(image, return_tensors="pt")
            out = self.model.generate(**inputs)
            caption = self.processor.decode(out[0], skip_special_tokens=True)
            return f"[image: {caption}]"
        except Exception:
            return ""

    def _get_structure_based_chunks(self, doc, pdf_path: str, body_font_size: float) -> List[Dict]:
        chunks = []
        current_topic = "General Content"
        current_text = []
        current_topic_page = 1

        header_threshold = body_font_size + 1.2

        for page_num, page in enumerate(doc):
            blocks = page.get_text("dict")["blocks"]
            blocks.sort(key=lambda b: b["bbox"][1])

            for b in blocks:
                if b["type"] == 1:
                    caption = self._generate_caption(b["image"])
                    if caption:
                        current_text.append(caption)
                    continue

                if "lines" in b:
                    for line in b["lines"]:
                        line_text = "".join([s["text"] for s in line["spans"]])
                        clean_line = self._clean_text(line_text)
                        if not clean_line:
                            continue

                        line_max_size = max([s["size"] for s in line["spans"]])

                        is_header = (
                            (line_max_size > header_threshold) and
                            (len(clean_line.split()) < 10) and
                            (not clean_line.strip().startswith(('"', "\u201c", "\u201d", "'"))) and
                            (clean_line[0].isupper() or clean_line[0].isdigit())
                        )

                        if is_header:
                            if current_text:
                                chunks.append({
                                    "doc_id": os.path.basename(pdf_path),
                                    "source_name": os.path.basename(pdf_path),
                                    "page": current_topic_page,
                                    "topic": current_topic,
                                    "chunk_id": len(chunks),
                                    "text": " ".join(current_text)
                                })
                                current_text = []

                            current_topic = clean_line
                            current_topic_page = page_num + 1
                        else:
                            current_text.append(clean_line)

        if current_text:
            chunks.append({
                "doc_id": os.path.basename(pdf_path),
                "source_name": os.path.basename(pdf_path),
                "page": current_topic_page,
                "topic": current_topic,
                "chunk_id": len(chunks),
                "text": " ".join(current_text)
            })

        return chunks

    def _get_paragraph_chunks(self, doc, pdf_path: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict]:
        full_text = ""
        for page in doc:
            full_text += page.get_text() + "\n\n"

        paragraphs = full_text.split('\n\n')
        chunks = []
        current_chunk = ""

        for para in paragraphs:
            clean_para = self._clean_text(para)
            if not clean_para:
                continue

            if len(current_chunk) + len(clean_para) > chunk_size:
                chunks.append({
                    "doc_id": os.path.basename(pdf_path),
                    "source_name": os.path.basename(pdf_path),
                    "page": "N/A",
                    "topic": "Paragraph Chunk",
                    "chunk_id": len(chunks),
                    "text": current_chunk.strip()
                })
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + " " + clean_para
            else:
                current_chunk += " " + clean_para

        if current_chunk:
            chunks.append({
                "doc_id": os.path.basename(pdf_path),
                "source_name": os.path.basename(pdf_path),
                "page": "N/A",
                "topic": "Paragraph Chunk",
                "chunk_id": len(chunks),
                "text": current_chunk.strip()
            })

        return chunks

    def process_pptx(self, pptx_path: str) -> List[Dict]:
        if not HAS_PPTX:
            raise ImportError("python-pptx is required to process .pptx files. Install with: pip install python-pptx")

        prs = Presentation(pptx_path)
        chunks = []
        fname = os.path.basename(pptx_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_title = "Untitled Slide"
            slide_texts = []

            for shape in slide.shapes:
              
                if shape.has_text_frame:
                    if shape == slide.shapes.title and shape.text.strip():
                        slide_title = self._clean_text(shape.text)
                    else:
                        text = self._clean_text(shape.text)
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            self._clean_text(cell.text) for cell in row.cells
                        )
                        if row_text.replace("|", "").strip():
                            slide_texts.append(row_text)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        caption = self._generate_caption(image_bytes)
                        if caption:
                            slide_texts.append(caption)
                    except Exception:
                        pass

            if slide_texts:
                chunks.append({
                    "doc_id": fname,
                    "source_name": fname,
                    "page": slide_num,
                    "topic": slide_title,
                    "chunk_id": len(chunks),
                    "text": " ".join(slide_texts)
                })

        return chunks

    def process_docx(self, docx_path: str) -> List[Dict]:
        if not HAS_DOCX:
            raise ImportError("python-docx is required to process .docx files. Install with: pip install python-docx")

        doc = DocxDocument(docx_path)
        fname = os.path.basename(docx_path)
        chunks = []
        current_topic = "General Content"
        current_text = []

        image_captions = self._extract_docx_images(doc)

        img_idx = 0
        for para in doc.paragraphs:
            text = self._clean_text(para.text)
            if not text:
                continue

            if para.style and para.style.name and para.style.name.startswith('Heading'):
                if current_text:
                    chunks.append({
                        "doc_id": fname,
                        "source_name": fname,
                        "page": "N/A",
                        "topic": current_topic,
                        "chunk_id": len(chunks),
                        "text": " ".join(current_text)
                    })
                    current_text = []
                current_topic = text
            else:
                current_text.append(text)

        if current_text:
            chunks.append({
                "doc_id": fname,
                "source_name": fname,
                "page": "N/A",
                "topic": current_topic,
                "chunk_id": len(chunks),
                "text": " ".join(current_text)
            })

        if image_captions:
            chunks.append({
                "doc_id": fname,
                "source_name": fname,
                "page": "N/A",
                "topic": "Document Images",
                "chunk_id": len(chunks),
                "text": " ".join(image_captions)
            })

        if len(chunks) < 3:
            chunks = self._docx_paragraph_chunks(doc, docx_path)

        return chunks

    def _extract_docx_images(self, doc) -> List[str]:
        captions = []
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    image_bytes = rel.target_part.blob
                    caption = self._generate_caption(image_bytes)
                    if caption:
                        captions.append(caption)
        except Exception:
            pass
        return captions

    def _docx_paragraph_chunks(self, doc, docx_path: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict]:
        full_text = "\n\n".join(p.text for p in doc.paragraphs)
        paragraphs = full_text.split('\n\n')
        chunks = []
        current_chunk = ""
        fname = os.path.basename(docx_path)

        for para in paragraphs:
            clean_para = self._clean_text(para)
            if not clean_para:
                continue

            if len(current_chunk) + len(clean_para) > chunk_size:
                chunks.append({
                    "doc_id": fname,
                    "source_name": fname,
                    "page": "N/A",
                    "topic": "Paragraph Chunk",
                    "chunk_id": len(chunks),
                    "text": current_chunk.strip()
                })
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + " " + clean_para
            else:
                current_chunk += " " + clean_para

        if current_chunk:
            chunks.append({
                "doc_id": fname,
                "source_name": fname,
                "page": "N/A",
                "topic": "Paragraph Chunk",
                "chunk_id": len(chunks),
                "text": current_chunk.strip()
            })

        return chunks

    def process_text_file(self, file_path: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict]:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            full_text = f.read()

        fname = os.path.basename(file_path)
        paragraphs = re.split(r'\n{2,}', full_text)
        chunks = []
        current_chunk = ""

        for para in paragraphs:
            clean_para = self._clean_text(para)
            if not clean_para:
                continue

            if len(current_chunk) + len(clean_para) > chunk_size:
                chunks.append({
                    "doc_id": fname,
                    "source_name": fname,
                    "page": "N/A",
                    "topic": "Text Chunk",
                    "chunk_id": len(chunks),
                    "text": current_chunk.strip()
                })
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + " " + clean_para
            else:
                current_chunk += " " + clean_para

        if current_chunk:
            chunks.append({
                "doc_id": fname,
                "source_name": fname,
                "page": "N/A",
                "topic": "Text Chunk",
                "chunk_id": len(chunks),
                "text": current_chunk.strip()
            })

        return chunks


# backward compatibility alias
PDFProcessor = DocumentProcessor


class TextPreprocessor:
    @staticmethod
    def clean_text(text: str) -> str:
        return re.sub(r'\s+', ' ', text).strip()

    @staticmethod
    def preprocess_chunks(chunks: List[Dict]) -> List[Dict]:
        processed_chunks = []
        for chunk in chunks:
            new_chunk = chunk.copy()
            new_chunk["text"] = TextPreprocessor.clean_text(chunk["text"])
            new_chunk["combined_text"] = f"{chunk['topic']}: {chunk['text']}"
            processed_chunks.append(new_chunk)
        return processed_chunks
